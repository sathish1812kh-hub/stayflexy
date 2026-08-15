import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG = RGBColor(11, 15, 25)         # #0b0f19
    COLOR_CARD = RGBColor(30, 41, 59)       # #1e293b
    COLOR_CYAN = RGBColor(56, 189, 248)     # #38bdf8
    COLOR_PURPLE = RGBColor(168, 85, 247)   # #a855f7
    COLOR_EMERALD = RGBColor(16, 185, 129)  # #10b981
    COLOR_AMBER = RGBColor(245, 158, 11)    # #f59e0b
    COLOR_WHITE = RGBColor(248, 250, 252)    # #f8fafc
    COLOR_MUTED = RGBColor(148, 163, 184)   # #94a3b8

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, category_text):
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(9.2), Inches(0.8))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN

        badge_box = slide.shapes.add_textbox(Inches(9.8), Inches(0.4), Inches(2.7), Inches(0.5))
        tf_b = badge_box.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = category_text.upper()
        p_b.alignment = PP_ALIGN.RIGHT
        p_b.font.name = "Arial"
        p_b.font.size = Pt(11)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_PURPLE

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "STAYFLEXI AUTONOMOUS ENGINEERING PLATFORM"
    p.font.name = "Arial"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = "Structural Phases (0-12), 4-Layer Intelligence Matrix & Platform v6.9.0 (v9)"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_WHITE
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "12 Microservices  |  4-Graph Intelligence Layer  |  v6.9.0 Certified Production Release"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_PURPLE

    # -------------------------------------------------------------
    # SLIDE 2: AI Pre-Task Operating Protocol
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "AI Pre-Task Operating Protocol (Phase 0 Startup)", "AI Startup Protocol")

    card_boot = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.3))
    card_boot.fill.solid()
    card_boot.fill.fore_color.rgb = COLOR_CARD
    tf_b = card_boot.text_frame
    tf_b.word_wrap = True

    p = tf_b.paragraphs[0]
    p.text = "📋 MANDATORY PHASE 0 BOOT SEQUENCE (Executed Before Any Code Change)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(12)

    steps_text = (
        "1. Read State File: Inspect docs/discovery/current-state.md (Active Task TSK-00132, commit hash).\n"
        "2. Verify Git Status: Run git status --porcelain to detect unrecorded drift or modified files.\n"
        "3. Verify Active Services: Probe ports 5432 (Postgres), 7687 (Neo4j), Codegraph SQLite. Auto-start if down.\n"
        "4. Inspect Active Tasks: Open docs/discovery/active-tasks.md to verify dependency conditions.\n"
        "5. Confirm Rulebook: Reference rules in docs/discovery/V5.2-Orchestrator.md for source-of-truth hierarchy.\n"
        "6. Summarize & Prompt: Output 'Context Recovery Report' and request confirmation prior to modifications.\n"
        "7. Codegraph Exploration First: Query codegraph_explore to inspect callers, blast radius & tests before edits."
    )
    p_s = tf_b.add_paragraph()
    p_s.text = steps_text
    p_s.font.size = Pt(12)
    p_s.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 3: Stateless Recovery Across Logouts
    # -------------------------------------------------------------
    slide3_new = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3_new)
    add_header(slide3_new, "Stateless Recovery: Context Preservation Across Logouts", "Stateless Architecture")

    card_stateless = slide3_new.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.3))
    card_stateless.fill.solid()
    card_stateless.fill.fore_color.rgb = COLOR_CARD
    tf_st = card_stateless.text_frame
    tf_st.word_wrap = True

    p = tf_st.paragraphs[0]
    p.text = "🔑 SCENARIO: Developer Logs Out & Logs In to a Fresh Conversation (0 Chat History)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(12)

    st_text = (
        "When a developer starts a completely new session, prompt transcript history is 100% erased.\n"
        "The system DOES NOT rely on chat history. Context is rebuilt deterministically in seconds:\n\n"
        "• Double-Buffered Markdown State: Phase 0 reads 4 git-tracked recovery files (current-state.md, active-tasks.md, etc.).\n"
        "• Git Commit Anchors: Current Git commit hash (56a6ed2...) anchors exact code state.\n"
        "• Neo4j Topology Sync: Native Neo4j (Port 7687) stores 83 nodes & 194 relationships across 12 microservices.\n"
        "• Codegraph AST Engine: SQLite database (.codegraph/codegraph.db) provides instant symbol & caller lookups.\n"
        "• Graphiti Decision Rationale: Graphiti loads memory packs JIT to restore ADR design rationale.\n\n"
        "Result: 100% deterministic session recovery without depending on previous chat logs!"
    )
    p_st_body = tf_st.add_paragraph()
    p_st_body.text = st_text
    p_st_body.font.size = Pt(12)
    p_st_body.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 4: Multi-Layer Context Verification Protocol
    # -------------------------------------------------------------
    slide4_new = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4_new)
    add_header(slide4_new, "Multi-Layer Context Verification Protocol (6 Independent Layers)", "Redundant Verification")

    rows_cnt = 7
    table_shape = slide4_new.shapes.add_table(rows_cnt, 4, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
    table = table_shape.table
    table.columns[0].width = Inches(1.3)
    table.columns[1].width = Inches(2.7)
    table.columns[2].width = Inches(4.5)
    table.columns[3].width = Inches(3.2)

    headers = ["Layer #", "Verification Method", "Operational Mechanism", "Failure Defense & Self-Healing"]
    for c_idx, h_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN

    v_rows = [
        ("Layer 1", "Double-Buffered Markdown State", "Reads current-state.md, project-context.md, active-tasks.md, decision-log.md.", "Detects state mismatch if 1 file is manually corrupted."),
        ("Layer 2", "Git Commit & Porcelain Shield", "Compares commit hash against git rev-parse HEAD & git status --porcelain.", "Detects uncommitted local changes or branch drift immediately."),
        ("Layer 3", "Neo4j Node Freshness Check", "Validates 83 nodes, 194 relationships, and Git commit linkage.", "Auto-starts native server via scripts/start-neo4j.ps1 on port 7687."),
        ("Layer 4", "Codegraph AST Blast-Radius", "Queries .codegraph/codegraph.db for callers and covering unit tests.", "Guards shared packages from broken callers before edits are made."),
        ("Layer 5", "Graphiti Temporal Rationale", "Queries Graphiti memory for active ADR decisions linked to active task.", "Prevents AI from undoing intentional architectural trade-offs."),
        ("Layer 6", "Service Health & Schema Probe", "Probes ports 5432, 7687, GraphQL Federation, and fast sync gate.", "Validates live schema and executes sync-graph.ps1 in <8s.")
    ]

    for r_idx, r_data in enumerate(v_rows):
        for c_idx, cell_value in enumerate(r_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_WHITE if c_idx == 0 else COLOR_MUTED
            if c_idx == 0:
                p.font.bold = True

    # -------------------------------------------------------------
    # SLIDE 5: Git Commit Anchors & State File Registry
    # -------------------------------------------------------------
    slide5_new = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5_new)
    add_header(slide5_new, "Git Commit Anchors & State File Registry", "Git Anchors")

    card_anchors = slide5_new.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.3))
    card_anchors.fill.solid()
    card_anchors.fill.fore_color.rgb = COLOR_CARD
    tf_an = card_anchors.text_frame
    tf_an.word_wrap = True

    p = tf_an.paragraphs[0]
    p.text = "⚓ IMMUTABLE GIT ANCHORS & STATE FILE REGISTRY"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD
    p.space_after = Pt(12)

    anchor_text = (
        "• Certified Release Tag: v6.9.0 (Release v6.9.0 - Complete Certified Stayflexi Platform)\n"
        "• Active Commit Anchor: 56a6ed2 (feat: embed Codegraph into specific lifecycle phases of V5.2 Orchestrator)\n"
        "• Parent Commit Anchor: 6022920 (feat: integrate Codegraph into startup protocol and compliance gate)\n"
        "• Grandparent Commit Anchor: e18dc6c (feat: implement Turbo CI caching, fast pre-push diffing, and dev commands)\n\n"
        "📂 Double-Buffered State Recovery Files:\n"
        "  1. docs/discovery/current-state.md — Active Task TSK-00132 & Git Commit Hash Anchor\n"
        "  2. docs/discovery/active-tasks.md — Pending Sprint Task Queue & Dependencies\n"
        "  3. docs/discovery/V5.2-Orchestrator.md — 13 System Phases & Multi-Graph Rulebook\n"
        "  4. GEMINI.md — Mandatory Pre-Task Startup Rules & Codegraph Protocol\n\n"
        "⚡ Verification Command: Phase 0 runs 'git rev-parse HEAD; git status --porcelain; powershell -File scripts/sync-graph.ps1'"
    )
    p_an_b = tf_an.add_paragraph()
    p_an_b.text = anchor_text
    p_an_b.font.size = Pt(12)
    p_an_b.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 6: Context Preservation & Git Backup
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Context Preservation & Git Backup Mechanism", "Context Persistence")

    s1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    s1.fill.solid()
    s1.fill.fore_color.rgb = COLOR_CARD
    tf1 = s1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "🔒 Stateless Session Recovery"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(10)

    p_body1 = (
        "• Chat history is NOT source of truth.\n\n"
        "• Double-Buffered State:\n"
        "  Progress recorded in current-state.md, project-context.md, active-tasks.md, decision-log.md.\n\n"
        "• JIT Awareness Restoration:\n"
        "  LLM rebuilds full state deterministically on boot."
    )
    p_b1 = tf1.add_paragraph()
    p_b1.text = p_body1
    p_b1.font.size = Pt(13)
    p_b1.font.color.rgb = COLOR_MUTED

    s2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2))
    s2.fill.solid()
    s2.fill.fore_color.rgb = COLOR_CARD
    tf2 = s2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "🌿 Git as Context Backup & Drift Shield"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE
    p.space_after = Pt(10)

    p_body2 = (
        "• Commit Anchor:\n"
        "  Session state files embed exact Git commit hash (2d4a0f674f...).\n\n"
        "• Freshness Scoring:\n"
        "  Neo4j compares node extractedAt vs Git LastModified commit date; sets isStale = true if drift detected.\n\n"
        "• Porcelain Check:\n"
        "  git status checks detect uncommitted local changes."
    )
    p_b2 = tf2.add_paragraph()
    p_b2.text = p_body2
    p_b2.font.size = Pt(13)
    p_b2.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 7: JIT Dynamic Phase Loading
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Just-In-Time (JIT) Dynamic Phase Loading", "Memory Architecture")

    card_jit = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    card_jit.fill.solid()
    card_jit.fill.fore_color.rgb = COLOR_CARD
    tf_j = card_jit.text_frame
    tf_j.word_wrap = True

    p = tf_j.paragraphs[0]
    p.text = "⚡ SELECTIVE JIT PHASE TRIGGER MATRIX"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD
    p.space_after = Pt(12)

    matrix_text = (
        "• Phases 0–3 (Core Boot & Discovery): Always loaded on initial session startup.\n"
        "• Phase 4–5 (Graphiti Memory & GraphQL): Loaded JIT when ADR decisions or API subgraphs change.\n"
        "• Phase 6 (Browser Intelligence): Loaded JIT ONLY when UI components (src/app/) or user journeys edit.\n"
        "• Phases 7–9 (Impact, Telemetry, Risk): Loaded JIT when Composite Risk Score (CRS) ≥ 30.\n"
        "• Phases 10–12 (Governance & Sync): Loaded JIT at task completion sign-off.\n\n"
        "💡 Realtime Example: Editing a backend controller triggers Phase 7 Impact Analysis (Cypher blast radius) "
        "but skips Phase 6 Browser Intelligence, saving 40,000 LLM context tokens!"
    )
    p_m = tf_j.add_paragraph()
    p_m.text = matrix_text
    p_m.font.size = Pt(13)
    p_m.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDES 8–12: 4-Layer Intelligence Tool Deep Dives
    # -------------------------------------------------------------
    graph_intelligence_tools = [
        ("Codegraph Deep Dive: Role, Use Case, Why & How", "Intelligence Layer (1/5)", "⚡ Codegraph (AST Code & Blast Radius)",
         "ROLE: Sub-Millisecond Code Symbol, AST Call Path & Blast-Radius Engine.\n\n"
         "PRIMARY USE CASE: Instant lookup of functions, classes, callers, dynamic dispatch, and covering unit test files.\n\n"
         "WHY WE USE IT: Reading files manually wastes tokens and risks blind refactoring. Codegraph gives callers before editing.\n\n"
         "HOW WE USE IT: Real-time file-watching daemon updates 48MB SQLite index (.codegraph/codegraph.db) in ~1s -> codegraph_explore.",
         "Query: codegraph_explore 'PricingCalculator computeNightlyRate' -> Returns verbatim source, 3 callers in ComputeDynamicRate.ts, and linked unit test suite (PricingCalculator.test.ts) in 30ms.", COLOR_CYAN),

        ("Graphify Deep Dive: Role, Use Case, Why & How", "Intelligence Layer (2/5)", "🔮 Graphify (Monorepo Corpus Analysis)",
         "ROLE: Codebase Knowledge Extractor, AST Analyzer & GraphRAG Engine.\n\n"
         "PRIMARY USE CASE: Converts 1,651 raw files (617k words) into structured AST communities & god-node reports.\n\n"
         "WHY WE USE IT: Monorepo modularity audits require global file clustering and circular dependency detection.\n\n"
         "HOW WE USE IT: Tree-sitter AST parsing -> Louvain community clustering -> Ontology Studio -> BFS/DFS GraphRAG search.",
         "Parsed C:\\Stayflexi: 9,154 nodes, 16,835 edges. Identified God Nodes (handleRouteError - 187 edges, Logger). Querying graphify query 'auth flow' returns multi-hop caller paths across microservices.", COLOR_PURPLE),
        
        ("Neo4j Deep Dive: Role, Use Case, Why & How", "Intelligence Layer (3/5)", "🌐 Neo4j Graph Database Engine",
         "ROLE: Structural Reality & Network Topology Engine (12 Services & Sagas).\n\n"
         "PRIMARY USE CASE: Serves as the structural source of truth for 12 microservices, 27 API routes, 10 Kafka topics, and releases.\n\n"
         "WHY WE USE IT: Relational DBs cannot query multi-hop dependency chains efficiently to calculate architecture blast radius.\n\n"
         "HOW WE USE IT: Native Zero-Docker Community 5.26 (Port 7687, Zulu JDK 21) with Cypher queries & MCP integration.",
         "Executing Cypher: MATCH (s:Microservice)-[:PUBLISHES_EVENT]->(t:KafkaTopic)<-[:SUBSCRIBES_EVENT]-(c:Microservice) RETURN s, t, c. Verifies complete saga choreography for booking creation.", COLOR_EMERALD),
        
        ("Graphiti Deep Dive: Role, Use Case, Why & How", "Intelligence Layer (4/5)", "🧠 Graphiti Temporal Memory",
         "ROLE: Dynamic Temporal Episodic Memory Engine (Context Memory Pack).\n\n"
         "PRIMARY USE CASE: Stores long-term decision rationale, ADRs, historical trade-offs, and feature evolution across session restarts.\n\n"
         "WHY WE USE IT: LLM chat sessions are stateless; prompt history is erased when a session ends. Prevents repeating past mistakes.\n\n"
         "HOW WE USE IT: Serializes decision episodes into structured memory packs loaded JIT during Phase 0 session boot.",
         "Retrieves decision rationale for PR #284 (customerType field). Reminds LLM that customerType was intentionally made nullable to prevent migration downtime across 5,000 hotel tenants.", COLOR_AMBER)
    ]

    for title, cat, tool_title, tool_text, tool_ex, color in graph_intelligence_tools:
        slide = prs.slides.add_slide(blank_slide_layout)
        set_slide_background(slide)
        add_header(slide, title, cat)

        s1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
        s1.fill.solid()
        s1.fill.fore_color.rgb = COLOR_CARD
        tf1 = s1.text_frame
        tf1.word_wrap = True
        p = tf1.paragraphs[0]
        p.text = tool_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(10)

        p_b1 = tf1.add_paragraph()
        p_b1.text = tool_text
        p_b1.font.size = Pt(11)
        p_b1.font.color.rgb = COLOR_WHITE

        s2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2))
        s2.fill.solid()
        s2.fill.fore_color.rgb = COLOR_CARD
        tf2 = s2.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        p.text = "⚡ Realtime Operational Usage & Example"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(10)

        p_b2 = tf2.add_paragraph()
        p_b2.text = tool_ex
        p_b2.font.size = Pt(12)
        p_b2.font.color.rgb = COLOR_MUTED

    # Slide 12: 4-Layer Intelligence Matrix Comparison
    slide12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide12)
    add_header(slide12, "4-Layer Intelligence Matrix: Comparative Role Breakdown", "Intelligence Layer (5/5)")

    rows_cnt = 5
    table_shape = slide12.shapes.add_table(rows_cnt, 5, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
    table = table_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(2.3)
    table.columns[4].width = Inches(2.0)

    headers = ["Tool Name", "Core Architectural Role", "Primary Input Data", "Output Artifact", "Query Engine"]
    for c_idx, h_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN

    matrix_rows = [
        ("Codegraph", "Sub-Millisecond Code AST & Blast Radius", "TypeScript AST symbols", ".codegraph/codegraph.db (48MB)", "codegraph_explore MCP"),
        ("Graphify", "Codebase Knowledge Extractor & GraphRAG", "1,651 Monorepo files", "Ontology Studio + JSON", "BFS / DFS GraphRAG CLI"),
        ("Neo4j", "Structural Topology & Distributed Sagas", "12 Services, APIs, Kafka", "Native Property Graph (Port 7687)", "Cypher Query Language"),
        ("Graphiti", "Temporal Episodic Decision Memory Engine", "ADRs, PR sign-offs", "Serialized Memory Pack", "Temporal Search API")
    ]

    for r_idx, r_data in enumerate(matrix_rows):
        for c_idx, cell_value in enumerate(r_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE if c_idx == 0 else COLOR_MUTED
            if c_idx == 0:
                p.font.bold = True

    # -------------------------------------------------------------
    # SLIDES 12–24: 4-Card Structural Deep Dives (Phases 0 to 12)
    # -------------------------------------------------------------
    # -------------------------------------------------------------
    # SLIDES 13–25: 4-Card Structural Deep Dives (Phases 0 to 12)
    # -------------------------------------------------------------
    structural_phase_slides = [
        ("Phase 0 Structural Deep Dive: Context Reconstruction", "System Phases (1/13)",
         "💡 WHY THIS PHASE EXISTS:\nChat history is erased across logouts; Phase 0 rebuilds exact operational context from codebase state to eliminate Artificial Intelligence (AI) hallucination.",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Reads 4 double-buffered recovery state files.\n2. Compares Git commit anchor vs git rev-parse HEAD.\n3. Runs git status --porcelain for code drift.\n4. Probes ports 5432, 7687, and Codegraph SQLite database.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nGit CLI (git status --porcelain), Codegraph AST Checker, Neo4j Bolt Client, File Inspector (view_file).",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nSession boots for TSK-00132. Reads current-state.md, verifies commit 56a6ed2, probes Neo4j port 7687 & Codegraph SQLite, emits Context Recovery Report.", COLOR_CYAN),
        
        ("Phase 1 Structural Deep Dive: Discovery & AST Indexing", "System Phases (2/13)",
         "💡 WHY THIS PHASE EXISTS:\nMonorepo contains 12 services and 9 packages; Phase 1 prevents blind refactoring by indexing all files, AST call paths, endpoints, and DB schemas.",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Queries codegraph_explore for symbol discovery & dynamic flows.\n2. Scans monorepo directory trees & package manifests.\n3. Identifies active REST endpoints, Prisma schemas & Kafka topics.\n4. Writes technology stack inventory & Codegraph AST symbol map.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nCodegraph AST Index (codegraph_explore), Ripgrep (grep_search), File System Crawler (list_dir).",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nRuns codegraph_explore 'TokenService' across shared-auth, auth-service, and booking-service; returns exact callers and dynamic dispatch paths in 30ms.", COLOR_PURPLE),
        
        ("Phase 2 Structural Deep Dive: Knowledge & Domain Modeling", "System Phases (3/13)",
         "💡 WHY THIS PHASE EXISTS:\nRaw code text lacks domain semantic structure; Phase 2 converts source symbols into typed graph entities (Service, Endpoint, Repo, Table, Saga).",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Parses Abstract Syntax Trees (AST) into Codegraph & Tree-sitter.\n2. Extracts class definitions, method signatures, and caller hierarchies.\n3. Maps DB tables, foreign keys & GraphQL resolvers.\n4. Appends typed node definitions to entity catalogs.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nCodegraph AST Engine (.codegraph/codegraph.db), Tree-sitter Parser, Node Catalog Builder.",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nModels CustomerProfile AST struct, updateCustomer GraphQL resolver, and Customer Prisma table model; appends typed node definitions to docs/discovery/NODE_CATALOG.md.", COLOR_EMERALD),
        
        ("Phase 3 Structural Deep Dive: Neo4j Graph Foundation", "System Phases (4/13)",
         "💡 WHY THIS PHASE EXISTS:\nRelational DBs cannot query multi-hop dependency chains efficiently; Neo4j physical property graph enables instant distributed system Topology Traversal.",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Ingests 83 typed topology nodes & 194 directed edges into Neo4j.\n2. Enforces Cypher uniqueness constraints & release linkage.\n3. Maps 12 microservices, 10 Kafka topics & 27 API routes.\n4. Validates graph health & node freshness timestamps.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nNative Neo4j Community 5.26 (Port 7687, Zulu JDK 21), Cypher Query Runner, Neo4j MCP Client.",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nIngests 83 nodes and 194 relationships; connects CustomerProfile node to 14 endpoint nodes and Kafka topic stayflexi.booking.created in Neo4j database on port 7687.", COLOR_AMBER),
        
        ("Phase 4 Structural Deep Dive: Graphiti Temporal Rationale Memory", "System Phases (5/13)",
         "💡 WHY THIS PHASE EXISTS:\nStandard DBs store current code state but lose historical design trade-off rationale; Phase 4 records decision episodes so AI never undoes Architectural Decision Records (ADRs).",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Captures Architectural Decision Records (ADRs) & PR trade-offs.\n2. Serializes decision episodes into dynamic memory packs.\n3. Links episodes to target task IDs & AST nodes.\n4. Loads memory packs JIT during Phase 0 session boot.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nGraphiti Memory Engine API, Episode Serializer, Temporal Search API.",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nSerializes ADR episode explaining WHY customerType is made nullable (prevents DB downtime across 5,000 hotel tenants); saves memory pack linked to task TSK-00132.", COLOR_PURPLE),
        
        ("Phase 5 Structural Deep Dive: GraphQL Supergraph Synchronization", "System Phases (6/13)",
         "💡 WHY THIS PHASE EXISTS:\n12 microservices expose federated subgraphs; Phase 5 prevents breaking client queries by verifying code-first schema compilation & supergraph rules.",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Compiles Pothos code-first schemas across microservices.\n2. Validates @key, @provides & @requires directives.\n3. Runs Apollo Router schema composition for supergraph.\n4. Checks for breaking field omissions or type mismatches.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nPothos Schema Builder, Apollo Router CLI (Port 8080), GraphQL Schema Linter.",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nExtends Pothos Customer entity in auth-service; verifies federated Customer entity resolution cleanly builds across Apollo Router; updates supergraph.graphql.", COLOR_EMERALD),
        
        ("Phase 6 Structural Deep Dive: Browser & User Journey Intelligence", "System Phases (7/13)",
         "💡 WHY THIS PHASE EXISTS:\nBackend code edits can cause silent visual bugs or layout breakage on User Interface (UI); Phase 6 validates real browser rendering & Document Object Model (DOM) trees.",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Launches headless Playwright browser instance.\n2. Executes guest checkout user flows on /bookings.\n3. Captures DOM tree snapshots & visual element bounds.\n4. Performs pixel-level visual regression diff analysis.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nHeadless Playwright, Puppeteer Engine, Visual Regression Diff Engine.",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nPlaywright script executes guest checkout flow on /bookings; verifies corporate customerType badge renders correctly without visual layout distortion; saves DOM diff artifacts.", COLOR_CYAN),
        
        ("Phase 7 Structural Deep Dive: Codegraph Blast Radius & Impact Analysis", "System Phases (8/13)",
         "💡 WHY THIS PHASE EXISTS:\nModifying a shared package or auth helper can break dozens of endpoints; Phase 7 calculates exact symbol Blast Radius (BR) and links covering test suites BEFORE writing code.",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Executes codegraph_explore on target symbols to find all callers.\n2. Traverses multi-hop Cypher relationships across 12 microservices in Neo4j.\n3. Auto-detects covering unit test files for the target symbols.\n4. Generates Blast-Radius Impact Report & test target list.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nCodegraph AST Engine (codegraph_explore), Neo4j Cypher Traversal Engine, Graphify Blast Assessor.",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nCodegraph exploration traces modification to TokenService payload; finds 14 dependent callers and auto-links TokenService.test.ts & auth.e2e.test.ts for validation.", COLOR_AMBER),
        
        ("Phase 8 Structural Deep Dive: Telemetry & Runtime Intelligence", "System Phases (9/13)",
         "💡 WHY THIS PHASE EXISTS:\nStatic code analysis cannot detect runtime latency degradation or memory leaks; Phase 8 inspects live OpenTelemetry (OTEL) traces and Prometheus metrics.",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Scrapes Prometheus metrics from /metrics endpoints.\n2. Collects OpenTelemetry (OTEL) distributed trace spans.\n3. Evaluates p95/p99 query latencies and error rates.\n4. Feeds latency metrics into Performance Risk (PR) vector.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nOpenTelemetry (OTEL) SDK, Prometheus Scraper (/metrics), Jaeger Client.",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nScrapes /metrics endpoint; detects 45ms DB query latency on customer profile lookups; feeds latency metrics into Performance Risk (PR) vector for risk scoring.", COLOR_PURPLE),
        
        ("Phase 9 Structural Deep Dive: Consequence Simulation & Risk Scoring", "System Phases (10/13)",
         "💡 WHY THIS PHASE EXISTS:\nChanges must be objectively categorized by risk before execution; Phase 9 evaluates 5 risk vectors to calculate Composite Risk Score (CRS).",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Evaluates Topology Risk (TR), Blast Radius (BR), Security Risk (SR), Performance Risk (PR), Complexity Risk (CR).\n2. Computes formula: CRS = (0.25 TR + 0.25 BR + 0.20 SR + 0.15 PR + 0.15 CR) * 10.0.\n3. Categorizes change: LOW (<30), MEDIUM (30-59), HIGH (60-84), CRITICAL (≥85).\n4. Dictates required approval gates (human sign-off).",
         "🛠️ TOOLS USED & REALTIME USAGE:\nComposite Risk Score (CRS) Engine, High-Concurrency Race Simulator.",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nEvaluates auth payload and schema changes -> calculates Composite Risk Score (CRS = 68 HIGH RISK) -> flags requirement for explicit human architect sign-off before proceeding.", COLOR_CYAN),
        
        ("Phase 10 Structural Deep Dive: Governance & Coordinated Execution", "System Phases (11/13)",
         "💡 WHY THIS PHASE EXISTS:\nEnforces organizational compliance, input sanitization, and pre-commit shields to execute coordinated multi-package refactoring safely.",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Executes coordinated edits across packages with caller awareness.\n2. Runs targeted unit tests auto-linked by Codegraph.\n3. Enforces Zod payload sanitization at Express middleware.\n4. Runs Husky pre-commit hooks for linting & type checks.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nCodegraph Test Linker, Governance Rule Evaluator, Zod Schema Sanitizer, Husky Pre-commit Hooks.",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nExecutes coordinated edits across shared-auth and auth-service; runs TokenService.test.ts; enforces Zod string enum sanitization for customerType; passes Husky.", COLOR_EMERALD),
        
        ("Phase 11 Structural Deep Dive: Multi-Engine Synchronization & Invariants", "System Phases (12/13)",
         "💡 WHY THIS PHASE EXISTS:\nPrevents database drift across system layers; guarantees code edits are instantly mirrored in Neo4j graph, Codegraph AST, Graphiti memory & GraphQL schemas.",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Validates invariant: Neo4j ≡ Codebase ∧ Codegraph ≡ Codebase ∧ GraphQL ≡ Codebase.\n2. Re-indexes AST node records for modified controllers in Codegraph.\n3. Updates Neo4j microservice topology and Kafka event edges.\n4. Halts completion gates if graph drift is detected in sync-task.ps1.",
         "🛠️ TOOLS USED & REALTIME USAGE:\nscripts/sync-graph.ps1 (<8s gate), Codegraph AST Watcher, Neo4j Bolt Client (Port 7687).",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nRuns scripts/sync-graph.ps1; verifies Postgres, Neo4j, and Codegraph report ONLINE; confirms zero drift across 83 nodes and 194 relationships in <8s.", COLOR_PURPLE),
        
        ("Phase 12 Structural Deep Dive: Session Persistence & Release Sign-Off", "System Phases (13/13)",
         "💡 WHY THIS PHASE EXISTS:\nGuarantees 100% deterministic session recovery for future session restarts by persisting state snapshots and anchoring Git commit hashes.",
         "⚙️ WHAT HAPPENS IN DETAIL:\n1. Double-buffers updated task status into current-state.md.\n2. Embeds new Git commit hash anchor (56a6ed2...).\n3. Signs transaction record in decision-log.md.\n4. Updates release tag v6.9.0-certified; system ready for next boot!",
         "🛠️ TOOLS USED & REALTIME USAGE:\nDouble-Buffer File Serializer, Git Commit & Tag Manager, Transaction Signer.",
         "🎯 REAL-TIME TASK WALKTHROUGH (TSK-00132):\nWrites task completion snapshot for TSK-00132 to current-state.md; embeds Git commit hash 56a6ed2...; updates release tag v6.9.0-certified.", COLOR_CYAN)
    ]

    for title, cat, why_text, what_text, how_text, task_text, color in structural_phase_slides:
        slide = prs.slides.add_slide(blank_slide_layout)
        set_slide_background(slide)
        add_header(slide, title, cat)

        cards = [
            (why_text, Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.5), COLOR_CYAN),
            (what_text, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.5), COLOR_WHITE),
            (how_text, Inches(0.8), Inches(4.2), Inches(5.6), Inches(2.5), COLOR_PURPLE),
            (task_text, Inches(6.8), Inches(4.2), Inches(5.7), Inches(2.5), COLOR_EMERALD)
        ]

        for c_text, left, top, width, height, c_color in cards:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_CARD
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = c_text
            p.font.size = Pt(11)
            p.font.color.rgb = c_color

    # -------------------------------------------------------------
    # SLIDE 26: REALTIME TASK SCENARIO: Pass-Through Flow
    # -------------------------------------------------------------
    slide26 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide26)
    add_header(slide26, "REALTIME TASK SCENARIO: Pass-Through Lifecycle Flow", "Realtime Task Scenario")

    scen_box = slide26.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.3))
    scen_box.fill.solid()
    scen_box.fill.fore_color.rgb = COLOR_CARD
    tf_sc = scen_box.text_frame
    tf_sc.word_wrap = True

    p = tf_sc.paragraphs[0]
    p.text = "🎯 LIVE TASK SCENARIO: PR #284 — Exposing Corporate customerType Fields Across Services"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(10)

    scen_text = (
        "Developer Request: Expose customerType in corporate customer profiles across shared-auth, auth-service, and booking-service.\n\n"
        "1. ⚡ Codegraph Phase:\n"
        "   Runs codegraph_explore on TokenService -> returns 14 dependent callers and links TokenService.test.ts unit test suite in 30ms.\n\n"
        "2. 🔮 Graphify Phase:\n"
        "   AST scans code structure -> maps 14 dependent endpoints and 3 microservice communities.\n\n"
        "3. 🌐 Neo4j Phase:\n"
        "   Runs Cypher impact query -> validates microservice topology & Kafka topic stayflexi.booking.created -> validates Neo4j == Codebase.\n\n"
        "4. 🧠 Graphiti Phase:\n"
        "   Queries prior decision memory -> records new ADR episode for nullable customerType -> serializes memory pack so future chat sessions retain full rationale."
    )
    p_st = tf_sc.add_paragraph()
    p_st.text = scen_text
    p_st.font.size = Pt(12)
    p_st.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDES 26–27: Data Stack & QA Stack
    # -------------------------------------------------------------
    stack_tables = [
        ("Data & Messaging Stack (Postgres, Redis, Kafka, Prisma)", "Infrastructure Stack (1/2)", [
            ("PostgreSQL 16", "Relational DB", "Primary persistent storage for bookings, inventory, and ledgers.", "Logical schema segregation across 12 microservices."),
            ("Redis 7", "Cache & Lock Store", "Handles JWT session caching, rate limits, and room reservation locks.", "Lua script acquiring atomic lock on lock:room:{id}."),
            ("Apache Kafka 7.5", "Event Bus", "Async messaging bus handling booking events & DLQ retries.", "Publishing BOOKING_CREATED to trigger notification workers."),
            ("Prisma 6.8 ORM", "Data Access Layer", "Consolidates 16 Prisma schemas into a single typed client.", "Generates strict TypeScript types for all repository queries.")
        ]),
        ("QA, Observability & Governance Stack", "Infrastructure Stack (2/2)", [
            ("Playwright & Puppeteer", "Browser Automation", "Automated user journey discovery, DOM snapshotting, visual diffs.", "Execution of guest checkout flows on /bookings."),
            ("OpenTelemetry", "Observability SDK", "Instruments HTTP/DB calls, exposes Prometheus metrics at /metrics.", "Scraping query latency metrics to calculate Performance Risk."),
            ("Zod", "Runtime Validation", "Input payload sanitization at Express middleware boundaries (wrapZod).", "Sanitizing incoming JSON payloads at API Gateway."),
            ("Jest & Supertest", "Automated Testing", "Unit test suites and HTTP endpoint contract integration tests.", "Running concurrency tests to verify lock race behavior."),
            ("Husky & Pre-commit", "Git Governance", "Pre-commit hooks for linting, type checks, and schema sync.", "Blocking uncommitted git commits prior to local commit.")
        ])
    ]

    for title, cat, rows in stack_tables:
        slide = prs.slides.add_slide(blank_slide_layout)
        set_slide_background(slide)
        add_header(slide, title, cat)

        rows_cnt = len(rows) + 1
        table_shape = slide.shapes.add_table(rows_cnt, 4, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
        table = table_shape.table
        table.columns[0].width = Inches(2.2)
        table.columns[1].width = Inches(2.3)
        table.columns[2].width = Inches(4.2)
        table.columns[3].width = Inches(3.0)

        headers = ["Tool / Technology", "Category", "Specific Use Case in Stayflexi", "Realtime Usage Example"]
        for c_idx, h_text in enumerate(headers):
            cell = table.cell(0, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            p.text = h_text
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN

        for r_idx, r_data in enumerate(rows):
            for c_idx, cell_value in enumerate(r_data):
                cell = table.cell(r_idx + 1, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
                p = cell.text_frame.paragraphs[0]
                p.text = cell_value
                p.font.size = Pt(11)
                p.font.color.rgb = COLOR_WHITE if c_idx == 0 else COLOR_MUTED
                if c_idx == 0:
                    p.font.bold = True

    # -------------------------------------------------------------
    # SLIDE 29: Summary & Deliverables Overview
    # -------------------------------------------------------------
    slide29 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide29)
    add_header(slide29, "Summary & Deliverables Overview", "Resources")

    s29 = slide29.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    s29.fill.solid()
    s29.fill.fore_color.rgb = COLOR_CARD
    tf29 = s29.text_frame
    tf29.word_wrap = True

    p = tf29.paragraphs[0]
    p.text = "🎉 STAYFLEXI AUTONOMOUS ENGINEERING DELIVERABLES"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(12)

    deliv_text = (
        "1. Interactive HTML Presentation Web App:\n"
        "   Path: stayflexi_engineering_presentation.html\n"
        "   Single-file dark-mode web app featuring KaTeX formulas, 29-slide selector, and arrow navigation.\n\n"
        "2. Native PowerPoint Widescreen Deck (.pptx):\n"
        "   Path: stayflexi_engineering_presentation_v9.pptx\n"
        "   Native 29-slide PowerPoint deck formatted with dark-mode styling, structured data tables, and widescreen layout.\n\n"
        "Generated by Stayflexi Autonomous Orchestrator • Version v6.9.0-certified (4-Layer Intelligence Matrix)"
    )
    p_dt = tf29.add_paragraph()
    p_dt.text = deliv_text
    p_dt.font.size = Pt(13)
    p_dt.font.color.rgb = COLOR_WHITE

    output_path = os.path.join("C:\\Stayflexi", "stayflexi_engineering_presentation_v9.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

    main_path = os.path.join("C:\\Stayflexi", "stayflexi_engineering_presentation.pptx")
    try:
        prs.save(main_path)
        print(f"Also updated main path: {main_path}")
    except Exception as e:
        print(f"Main path locked by viewer ({e}); generated v9 file at: {output_path}")

if __name__ == "__main__":
    create_presentation()
